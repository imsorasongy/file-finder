import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import os
import sys
import time
import threading
import subprocess
import sqlite3
import datetime
from pathlib import Path


_CREATE_NO_WINDOW = 0x08000000 if os.name == 'nt' else 0
INDEX_DB_PATH = str(Path.home() / '.파일찾기_색인.db')


# ============================================================
# 텍스트 추출기
# ============================================================
TEXT_EXTENSIONS = {
    'txt', 'md', 'markdown', 'rst', 'log', 'csv', 'tsv',
    'py', 'pyw', 'js', 'mjs', 'ts', 'tsx', 'jsx',
    'html', 'htm', 'xml', 'json', 'yaml', 'yml', 'toml', 'ini', 'cfg', 'conf',
    'css', 'scss', 'sass',
    'c', 'cc', 'cpp', 'h', 'hpp', 'java', 'kt', 'go', 'rs', 'rb', 'php',
    'sh', 'bat', 'ps1',
    'sql', 'r',
    'vue', 'svelte',
}
BINARY_SKIP_EXTENSIONS = {
    'exe', 'dll', 'so', 'dylib', 'bin', 'o', 'obj', 'lib', 'a',
    'png', 'jpg', 'jpeg', 'gif', 'bmp', 'ico', 'webp', 'tiff', 'tif', 'svg',
    'mp3', 'wav', 'flac', 'ogg', 'm4a', 'aac',
    'mp4', 'avi', 'mov', 'mkv', 'wmv', 'webm', 'flv',
    'zip', 'rar', '7z', 'tar', 'gz', 'bz2', 'xz',
    'iso', 'dmg',
    'ttf', 'otf', 'woff', 'woff2',
    'class', 'pyc', 'pyo',
    'psd', 'ai', 'eps', 'sketch',
    'db', 'sqlite', 'sqlite3', 'mdb',
}
CONTENT_SUPPORTED_EXTS = TEXT_EXTENSIONS | {'pdf', 'hwp', 'docx', 'pptx', 'xlsx'}


def extract_text_from_file(filepath, max_bytes):
    ext = filepath.rsplit('.', 1)[-1].lower() if '.' in filepath else ''
    if ext in BINARY_SKIP_EXTENSIONS:
        return None
    try:
        fsize = os.path.getsize(filepath)
    except OSError:
        return None
    if fsize > max_bytes:
        return None

    if ext == 'pdf':
        try:
            from pypdf import PdfReader
            r = PdfReader(filepath)
            parts = []
            for p in r.pages:
                try:
                    parts.append(p.extract_text() or '')
                except Exception:
                    continue
            return '\n'.join(parts)
        except Exception:
            return None

    if ext == 'hwp':
        try:
            import olefile, zlib, re
            ole = olefile.OleFileIO(filepath)
            texts = []
            if ole.exists('PrvText'):
                try:
                    raw = ole.openstream('PrvText').read()
                    texts.append(raw.decode('utf-16-le', errors='ignore'))
                except Exception:
                    pass
            compressed = False
            if ole.exists('FileHeader'):
                try:
                    hdr = ole.openstream('FileHeader').read()
                    if len(hdr) >= 37:
                        compressed = bool(hdr[36] & 0x01)
                except Exception:
                    pass
            for entry in ole.listdir():
                if len(entry) >= 2 and entry[0] == 'BodyText' and entry[1].startswith('Section'):
                    try:
                        raw = ole.openstream(entry).read()
                        if compressed:
                            try:
                                raw = zlib.decompress(raw, -15)
                            except zlib.error:
                                continue
                        decoded = raw.decode('utf-16-le', errors='ignore')
                        cleaned = re.sub(
                            r'[^\uAC00-\uD7A3\u3131-\u318E'
                            r'a-zA-Z0-9'
                            r' .,!?()\[\]{}:;\-\n/%*+=\'"<>]',
                            ' ', decoded
                        )
                        cleaned = re.sub(r'\s{3,}', '  ', cleaned)
                        texts.append(cleaned)
                    except Exception:
                        continue
            ole.close()
            return '\n'.join(texts) if texts else None
        except Exception:
            return None

    if ext == 'docx':
        try:
            from docx import Document
            doc = Document(filepath)
            parts = [p.text for p in doc.paragraphs]
            for t in doc.tables:
                for row in t.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            return '\n'.join(parts)
        except Exception:
            return None

    if ext == 'pptx':
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for slide in prs.slides:
                for sh in slide.shapes:
                    if sh.has_text_frame:
                        parts.append(sh.text_frame.text)
                    if sh.has_table:
                        for row in sh.table.rows:
                            for cell in row.cells:
                                parts.append(cell.text)
            return '\n'.join(parts)
        except Exception:
            return None

    if ext == 'xlsx':
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, data_only=True, read_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    parts.append(' '.join(str(v) for v in row if v is not None))
            wb.close()
            return '\n'.join(parts)
        except Exception:
            return None

    if ext in TEXT_EXTENSIONS or ext == '':
        try:
            with open(filepath, 'rb') as f:
                raw = f.read()
            for enc in ('utf-8-sig', 'utf-8', 'cp949', 'euc-kr', 'latin-1'):
                try:
                    return raw.decode(enc)
                except UnicodeDecodeError:
                    continue
        except Exception:
            return None
    return None


def extract_preview_text(filepath, max_bytes):
    """빠른 미리보기용 — 앞쪽 페이지/섹션만 추출."""
    ext = filepath.rsplit('.', 1)[-1].lower() if '.' in filepath else ''
    if ext in BINARY_SKIP_EXTENSIONS:
        return None
    try:
        fsize = os.path.getsize(filepath)
    except OSError:
        return None
    if fsize > max_bytes:
        return None

    if ext == 'pdf':
        try:
            from pypdf import PdfReader
            r = PdfReader(filepath)
            parts = []
            for p in r.pages[:3]:
                try:
                    parts.append(p.extract_text() or '')
                except Exception:
                    continue
            return '\n'.join(parts)
        except Exception:
            return None

    if ext == 'hwp':
        # PrvText만 (이미 본문 요약, 빠름)
        try:
            import olefile
            ole = olefile.OleFileIO(filepath)
            txt = ''
            if ole.exists('PrvText'):
                try:
                    raw = ole.openstream('PrvText').read()
                    txt = raw.decode('utf-16-le', errors='ignore')
                except Exception:
                    pass
            ole.close()
            return txt if txt else None
        except Exception:
            return None

    if ext == 'docx':
        try:
            from docx import Document
            doc = Document(filepath)
            parts = []
            for i, p in enumerate(doc.paragraphs):
                if i >= 100:
                    break
                parts.append(p.text)
            return '\n'.join(parts)
        except Exception:
            return None

    if ext == 'pptx':
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for i, slide in enumerate(prs.slides):
                if i >= 5:
                    break
                for sh in slide.shapes:
                    if sh.has_text_frame:
                        parts.append(sh.text_frame.text)
            return '\n'.join(parts)
        except Exception:
            return None

    if ext == 'xlsx':
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, data_only=True, read_only=True)
            parts = []
            if wb.worksheets:
                sheet = wb.worksheets[0]
                for i, row in enumerate(sheet.iter_rows(values_only=True)):
                    if i >= 100:
                        break
                    parts.append(' '.join(str(v) for v in row if v is not None))
            wb.close()
            return '\n'.join(parts)
        except Exception:
            return None

    # 텍스트 — 앞 100KB만 읽음
    if ext in TEXT_EXTENSIONS or ext == '':
        try:
            with open(filepath, 'rb') as f:
                raw = f.read(100_000)
            for enc in ('utf-8-sig', 'utf-8', 'cp949', 'euc-kr', 'latin-1'):
                try:
                    return raw.decode(enc)
                except UnicodeDecodeError:
                    continue
        except Exception:
            return None
    return None


def render_pdf_first_page_as_image(filepath, zoom=2.0):
    """스캔 PDF 등 텍스트 없는 PDF를 위한 첫 페이지 이미지 렌더링. PIL.Image 반환."""
    try:
        import fitz  # PyMuPDF
        from PIL import Image
        import io
        doc = fitz.open(filepath)
        if len(doc) == 0:
            doc.close()
            return None
        page = doc[0]
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes("png")
        doc.close()
        img = Image.open(io.BytesIO(img_bytes))
        img.load()
        return img
    except Exception:
        return None


def find_snippet(text, query, case_sensitive, ctx_before=40, ctx_after=80):
    if not text or not query:
        return None
    haystack = text if case_sensitive else text.lower()
    needle = query if case_sensitive else query.lower()
    idx = haystack.find(needle)
    if idx < 0:
        return None
    start = max(0, idx - ctx_before)
    end = min(len(text), idx + len(query) + ctx_after)
    snippet = text[start:end]
    snippet = snippet.replace('\n', ' ').replace('\r', ' ').replace('\t', ' ')
    while '  ' in snippet:
        snippet = snippet.replace('  ', ' ')
    snippet = snippet.strip()
    if start > 0:
        snippet = '…' + snippet
    if end < len(text):
        snippet = snippet + '…'
    return snippet


# ============================================================
# 색인 관리자 (SQLite FTS5 + trigram)
# ============================================================
class IndexManager:
    def __init__(self, db_path=INDEX_DB_PATH):
        self.db_path = db_path
        self._init_db()

    def _connect(self):
        conn = sqlite3.connect(self.db_path, timeout=30)
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("PRAGMA synchronous=NORMAL")
        return conn

    def _init_db(self):
        conn = self._connect()
        conn.executescript('''
            CREATE TABLE IF NOT EXISTS indexed_folders (
                id INTEGER PRIMARY KEY,
                path TEXT UNIQUE NOT NULL,
                added_at REAL NOT NULL
            );
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY,
                path TEXT UNIQUE NOT NULL,
                name TEXT NOT NULL,
                ext TEXT,
                size INTEGER,
                mtime REAL,
                indexed_at REAL NOT NULL,
                has_content INTEGER NOT NULL DEFAULT 0
            );
            CREATE INDEX IF NOT EXISTS idx_files_ext ON files(ext);
            CREATE VIRTUAL TABLE IF NOT EXISTS files_fts USING fts5(
                name, content, tokenize='trigram'
            );
        ''')
        conn.commit()
        conn.close()

    def add_folder(self, path):
        path = os.path.normpath(path)
        if not os.path.isdir(path):
            raise ValueError(f"폴더가 존재하지 않습니다: {path}")
        conn = self._connect()
        try:
            conn.execute(
                "INSERT OR IGNORE INTO indexed_folders(path, added_at) VALUES(?, ?)",
                (path, time.time())
            )
            conn.commit()
        finally:
            conn.close()

    def remove_folder(self, path):
        path = os.path.normpath(path)
        conn = self._connect()
        try:
            conn.execute("DELETE FROM indexed_folders WHERE path = ?", (path,))
            # 해당 폴더의 파일들도 제거
            like = path.rstrip(os.sep) + os.sep + '%'
            rows = conn.execute(
                "SELECT id FROM files WHERE path = ? OR path LIKE ?",
                (path, like)
            ).fetchall()
            ids = [r[0] for r in rows]
            if ids:
                placeholders = ','.join('?' * len(ids))
                conn.execute(f"DELETE FROM files_fts WHERE rowid IN ({placeholders})", ids)
                conn.execute(f"DELETE FROM files WHERE id IN ({placeholders})", ids)
            conn.commit()
        finally:
            conn.close()

    def list_folders(self):
        conn = self._connect()
        try:
            rows = conn.execute(
                "SELECT path, added_at FROM indexed_folders ORDER BY path"
            ).fetchall()
            return rows
        finally:
            conn.close()

    def stats(self):
        conn = self._connect()
        try:
            total = conn.execute("SELECT COUNT(*) FROM files").fetchone()[0]
            with_content = conn.execute(
                "SELECT COUNT(*) FROM files WHERE has_content=1"
            ).fetchone()[0]
            last = conn.execute("SELECT MAX(indexed_at) FROM files").fetchone()[0]
        finally:
            conn.close()
        db_size = os.path.getsize(self.db_path) if os.path.exists(self.db_path) else 0
        return {
            'total_files': total,
            'with_content': with_content,
            'last_indexed': last,
            'db_size': db_size,
        }

    def build(self, incremental, max_bytes, skip_hidden,
              progress_cb, stop_flag):
        """incremental=True: mtime 비교해 신규/변경만. False: 전체 재색인."""
        folders = [r[0] for r in self.list_folders()]
        if not folders:
            return {'indexed': 0, 'skipped': 0, 'removed': 0, 'scanned': 0}

        conn = self._connect()
        try:
            # 전체 재색인 시 기존 데이터 삭제
            if not incremental:
                conn.execute("DELETE FROM files_fts")
                conn.execute("DELETE FROM files")
                conn.commit()

            # 증분 시: 색인 대상 외 폴더의 파일은 정리되지 않음 (제거는 remove_folder에서 처리)
            # 폴더 밖으로 이동/삭제된 파일 제거
            if incremental:
                alive_ids = set()

            # 기존 색인 정보 (증분용)
            existing = {}
            if incremental:
                for row in conn.execute("SELECT id, path, mtime FROM files"):
                    existing[row[1]] = (row[0], row[2])

            indexed = 0
            skipped = 0
            scanned = 0
            seen_paths = set()

            for folder in folders:
                if stop_flag.is_set():
                    break
                for root_dir, dirs, files in os.walk(folder, onerror=lambda e: None):
                    if stop_flag.is_set():
                        break
                    if skip_hidden:
                        dirs[:] = [d for d in dirs
                                   if not self._is_hidden(os.path.join(root_dir, d))]
                    for f in files:
                        if stop_flag.is_set():
                            break
                        scanned += 1
                        fpath = os.path.join(root_dir, f)
                        seen_paths.add(fpath)
                        try:
                            st = os.stat(fpath)
                        except OSError:
                            continue

                        ext = f.rsplit('.', 1)[-1].lower() if '.' in f else ''

                        if incremental and fpath in existing:
                            old_id, old_mtime = existing[fpath]
                            if abs(old_mtime - st.st_mtime) < 0.001:
                                skipped += 1
                                if scanned % 200 == 0:
                                    progress_cb(scanned, indexed, skipped, f)
                                continue
                            # 변경된 파일: 재색인
                            self._index_one(conn, fpath, f, ext, st, max_bytes,
                                             replace_id=old_id)
                        else:
                            self._index_one(conn, fpath, f, ext, st, max_bytes)
                        indexed += 1
                        if indexed % 50 == 0:
                            conn.commit()
                            progress_cb(scanned, indexed, skipped, f)

            conn.commit()

            # 증분: 스캔하지 못한(= 삭제된 것으로 간주되는) 파일 정리
            removed = 0
            if incremental and not stop_flag.is_set():
                # 색인 대상 폴더 하위 중 이번에 안 보인 파일만 삭제
                to_remove = []
                for path, (fid, _) in existing.items():
                    if path in seen_paths:
                        continue
                    # 해당 파일이 색인 대상 폴더 하위인지 확인
                    for folder in folders:
                        folder_norm = folder.rstrip(os.sep)
                        if path == folder_norm or path.startswith(folder_norm + os.sep):
                            to_remove.append(fid)
                            break
                if to_remove:
                    placeholders = ','.join('?' * len(to_remove))
                    conn.execute(
                        f"DELETE FROM files_fts WHERE rowid IN ({placeholders})",
                        to_remove
                    )
                    conn.execute(
                        f"DELETE FROM files WHERE id IN ({placeholders})",
                        to_remove
                    )
                    conn.commit()
                    removed = len(to_remove)

            progress_cb(scanned, indexed, skipped, '')
            return {'indexed': indexed, 'skipped': skipped,
                    'removed': removed, 'scanned': scanned}
        finally:
            conn.close()

    def _index_one(self, conn, fpath, fname, ext, st, max_bytes, replace_id=None):
        content = extract_text_from_file(fpath, max_bytes) or ''
        has_content = 1 if content else 0
        now = time.time()
        if replace_id is not None:
            conn.execute(
                "UPDATE files SET name=?, ext=?, size=?, mtime=?, "
                "indexed_at=?, has_content=? WHERE id=?",
                (fname, ext, st.st_size, st.st_mtime, now, has_content, replace_id)
            )
            conn.execute("DELETE FROM files_fts WHERE rowid = ?", (replace_id,))
            conn.execute(
                "INSERT INTO files_fts(rowid, name, content) VALUES(?, ?, ?)",
                (replace_id, fname, content)
            )
        else:
            cur = conn.execute(
                "INSERT OR REPLACE INTO files(path, name, ext, size, mtime, "
                "indexed_at, has_content) VALUES(?, ?, ?, ?, ?, ?, ?)",
                (fpath, fname, ext, st.st_size, st.st_mtime, now, has_content)
            )
            new_id = cur.lastrowid
            conn.execute(
                "INSERT OR REPLACE INTO files_fts(rowid, name, content) "
                "VALUES(?, ?, ?)",
                (new_id, fname, content)
            )

    @staticmethod
    def _is_hidden(path):
        name = os.path.basename(path)
        if name.startswith('.'):
            return True
        try:
            import ctypes
            attrs = ctypes.windll.kernel32.GetFileAttributesW(str(path))
            if attrs == -1:
                return False
            return bool(attrs & 2) or bool(attrs & 4)
        except Exception:
            return False

    def search(self, name_q, ext_q, content_q, case_sensitive, limit=5000):
        """색인에서 검색. 3자 미만은 LIKE 폴백."""
        conn = self._connect()
        try:
            where = []
            params = []
            joins = []

            fts_name_parts = []
            fts_content_parts = []

            # NAME
            if name_q:
                if len(name_q) >= 3:
                    fts_name_parts.append(name_q)
                else:
                    where.append("files.name LIKE ? COLLATE NOCASE" if not case_sensitive
                                 else "files.name LIKE ?")
                    params.append(f'%{name_q}%')

            # CONTENT
            if content_q:
                if len(content_q) >= 3:
                    fts_content_parts.append(content_q)
                else:
                    # LIKE 폴백 — FTS 테이블의 content 컬럼은 LIKE 가능
                    where.append("files_fts.content LIKE ? COLLATE NOCASE" if not case_sensitive
                                 else "files_fts.content LIKE ?")
                    params.append(f'%{content_q}%')

            # 확장자
            if ext_q:
                exts = [e.strip().lstrip('.').lower() for e in ext_q.split(',') if e.strip()]
                if exts:
                    ph = ','.join('?' * len(exts))
                    where.append(f"files.ext IN ({ph})")
                    params.extend(exts)

            # FTS MATCH 구성
            fts_where = []
            fts_params = []
            if fts_name_parts and fts_content_parts:
                # 두 컬럼 모두 FTS
                match_expr = (f'name : "{self._escape_fts(fts_name_parts[0])}" AND '
                              f'content : "{self._escape_fts(fts_content_parts[0])}"')
                fts_where.append("files_fts MATCH ?")
                fts_params.append(match_expr)
            elif fts_name_parts:
                fts_where.append("files_fts.name MATCH ?")
                fts_params.append(f'"{self._escape_fts(fts_name_parts[0])}"')
            elif fts_content_parts:
                fts_where.append("files_fts.content MATCH ?")
                fts_params.append(f'"{self._escape_fts(fts_content_parts[0])}"')

            # JOIN 필요 여부
            need_fts = bool(fts_where) or any('files_fts' in w for w in where)
            if need_fts:
                base = "FROM files_fts JOIN files ON files.id = files_fts.rowid"
            else:
                base = "FROM files"

            all_where = fts_where + where
            all_params = fts_params + params
            sql = "SELECT files.path, files.name, files.size, files.mtime"
            # snippet
            if fts_content_parts:
                sql += (", snippet(files_fts, 1, '', '', '…', 30)")
            else:
                sql += ", ''"
            sql += " " + base
            if all_where:
                sql += " WHERE " + " AND ".join(all_where)
            sql += f" LIMIT {int(limit)}"

            rows = conn.execute(sql, all_params).fetchall()
            return rows
        finally:
            conn.close()

    @staticmethod
    def _escape_fts(s):
        return s.replace('"', '""')


# ============================================================
# 색인 관리 대화상자
# ============================================================
class IndexDialog(tk.Toplevel):
    def __init__(self, parent, index_mgr, on_close=None):
        super().__init__(parent)
        self.title("색인 관리")
        self.geometry("760x560")
        self.transient(parent)
        self.index_mgr = index_mgr
        self.on_close = on_close
        self.stop_flag = threading.Event()
        self.worker = None

        self._build_ui()
        self._refresh()
        self.protocol("WM_DELETE_WINDOW", self._close)

    def _build_ui(self):
        pad = 12

        top = ttk.Frame(self, padding=(pad, pad, pad, 6))
        top.pack(fill='x')
        ttk.Label(top, text="색인 대상 폴더", font=('Segoe UI', 11, 'bold')).pack(side='left')

        body = ttk.Frame(self, padding=(pad, 0, pad, 0))
        body.pack(fill='both', expand=True)

        list_frame = ttk.Frame(body)
        list_frame.pack(side='left', fill='both', expand=True)
        self.folder_list = tk.Listbox(list_frame, height=8)
        self.folder_list.pack(side='left', fill='both', expand=True)
        sb = ttk.Scrollbar(list_frame, orient='vertical', command=self.folder_list.yview)
        sb.pack(side='right', fill='y')
        self.folder_list.configure(yscrollcommand=sb.set)

        btns = ttk.Frame(body, padding=(8, 0))
        btns.pack(side='left', fill='y')
        ttk.Button(btns, text="폴더 추가", command=self._add_folder, width=14).pack(pady=3)
        ttk.Button(btns, text="선택 삭제", command=self._remove_folder, width=14).pack(pady=3)

        # 통계
        stat_frame = ttk.LabelFrame(self, text=" 색인 상태 ", padding=pad)
        stat_frame.pack(fill='x', padx=pad, pady=(10, 4))
        self.stat_label = ttk.Label(stat_frame, text="", justify='left')
        self.stat_label.pack(anchor='w')

        # 옵션
        opt_frame = ttk.LabelFrame(self, text=" 색인 옵션 ", padding=pad)
        opt_frame.pack(fill='x', padx=pad, pady=4)
        op = ttk.Frame(opt_frame)
        op.pack(fill='x')
        ttk.Label(op, text="내용 추출 최대 파일 크기 (MB)").pack(side='left')
        self.maxsize_var = tk.StringVar(value="20")
        ttk.Entry(op, textvariable=self.maxsize_var, width=6).pack(side='left', padx=(6, 12))
        self.skip_hidden_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(op, text="숨김/시스템 폴더 제외",
                        variable=self.skip_hidden_var).pack(side='left')

        # 빌드 버튼
        build_frame = ttk.Frame(self, padding=(pad, 0, pad, 4))
        build_frame.pack(fill='x')
        self.incremental_btn = ttk.Button(
            build_frame, text="🔄  증분 업데이트",
            command=lambda: self._start_build(incremental=True)
        )
        self.incremental_btn.pack(side='left', padx=(0, 6))
        self.full_btn = ttk.Button(
            build_frame, text="🗑  전체 재색인",
            command=lambda: self._start_build(incremental=False)
        )
        self.full_btn.pack(side='left', padx=(0, 6))
        self.stop_btn = ttk.Button(
            build_frame, text="■  중지", command=self._stop_build, state='disabled'
        )
        self.stop_btn.pack(side='left')

        # 진행 상태
        prog_frame = ttk.Frame(self, padding=(pad, 4, pad, pad))
        prog_frame.pack(fill='x')
        self.progress_label = ttk.Label(prog_frame, text="대기 중", foreground='#666')
        self.progress_label.pack(anchor='w')
        self.progress = ttk.Progressbar(prog_frame, mode='indeterminate')
        self.progress.pack(fill='x', pady=(4, 0))

        # 닫기
        ttk.Button(self, text="닫기", command=self._close, width=12).pack(pady=(0, pad))

    def _refresh(self):
        self.folder_list.delete(0, 'end')
        for path, added_at in self.index_mgr.list_folders():
            dt = datetime.datetime.fromtimestamp(added_at).strftime('%Y-%m-%d')
            self.folder_list.insert('end', f"{path}    ({dt} 추가)")
        st = self.index_mgr.stats()
        size_mb = st['db_size'] / (1024 * 1024)
        last = (datetime.datetime.fromtimestamp(st['last_indexed']).strftime('%Y-%m-%d %H:%M')
                if st['last_indexed'] else '(없음)')
        self.stat_label.configure(text=(
            f"색인된 파일: {st['total_files']:,}개  "
            f"(내용 추출됨: {st['with_content']:,}개)\n"
            f"DB 크기: {size_mb:,.1f} MB\n"
            f"최근 색인: {last}"
        ))

    def _add_folder(self):
        d = filedialog.askdirectory(parent=self, title="색인할 폴더 선택")
        if not d:
            return
        try:
            self.index_mgr.add_folder(d)
        except Exception as e:
            messagebox.showerror("오류", str(e), parent=self)
            return
        self._refresh()

    def _remove_folder(self):
        sel = self.folder_list.curselection()
        if not sel:
            return
        item = self.folder_list.get(sel[0])
        path = item.split('    ')[0]
        if not messagebox.askyesno(
            "확인",
            f"다음 폴더를 색인에서 제거하시겠습니까?\n\n{path}\n\n"
            "(폴더의 색인 데이터도 함께 삭제됩니다.)",
            parent=self
        ):
            return
        try:
            self.index_mgr.remove_folder(path)
        except Exception as e:
            messagebox.showerror("오류", str(e), parent=self)
            return
        self._refresh()

    def _start_build(self, incremental):
        if self.worker and self.worker.is_alive():
            return
        if not self.index_mgr.list_folders():
            messagebox.showinfo("알림", "먼저 색인할 폴더를 추가하세요.", parent=self)
            return
        try:
            max_mb = float(self.maxsize_var.get())
        except ValueError:
            messagebox.showerror("오류", "최대 파일 크기는 숫자여야 합니다.", parent=self)
            return
        if not incremental:
            if not messagebox.askyesno(
                "확인",
                "전체 재색인을 시작합니다. 기존 색인 데이터가 모두 삭제되고 다시 만들어집니다.\n\n계속할까요?",
                parent=self
            ):
                return

        self.stop_flag.clear()
        self.incremental_btn.configure(state='disabled')
        self.full_btn.configure(state='disabled')
        self.stop_btn.configure(state='normal')
        self.progress.configure(mode='indeterminate')
        self.progress.start(10)
        self.progress_label.configure(
            text=f"{'증분 업데이트' if incremental else '전체 재색인'} 시작..."
        )

        max_bytes = int(max_mb * 1024 * 1024)
        skip_hidden = self.skip_hidden_var.get()

        def progress_cb(scanned, indexed, skipped, current_file):
            self.after(0, lambda: self.progress_label.configure(
                text=(f"스캔 {scanned:,} · 색인 {indexed:,} · 건너뜀 {skipped:,}"
                      + (f"  ·  {current_file[:40]}" if current_file else ''))
            ))

        def worker():
            try:
                result = self.index_mgr.build(
                    incremental=incremental, max_bytes=max_bytes,
                    skip_hidden=skip_hidden,
                    progress_cb=progress_cb, stop_flag=self.stop_flag
                )
                self.after(0, self._finish_build, result, None)
            except Exception as e:
                self.after(0, self._finish_build, None, str(e))

        self.worker = threading.Thread(target=worker, daemon=True)
        self.worker.start()

    def _stop_build(self):
        self.stop_flag.set()

    def _finish_build(self, result, error):
        self.progress.stop()
        self.progress.configure(mode='determinate', value=0)
        self.incremental_btn.configure(state='normal')
        self.full_btn.configure(state='normal')
        self.stop_btn.configure(state='disabled')
        if error:
            self.progress_label.configure(text=f"오류: {error}")
            messagebox.showerror("오류", f"색인 중 오류가 발생했습니다:\n{error}", parent=self)
        elif self.stop_flag.is_set():
            self.progress_label.configure(
                text=(f"중지됨 — 스캔 {result['scanned']:,} · 색인 {result['indexed']:,}"
                      if result else "중지됨")
            )
        else:
            self.progress_label.configure(
                text=(f"완료 — 스캔 {result['scanned']:,} · "
                      f"새 색인 {result['indexed']:,} · "
                      f"건너뜀 {result['skipped']:,} · "
                      f"삭제 {result['removed']:,}")
            )
        self._refresh()
        if self.on_close:
            self.on_close()

    def _close(self):
        if self.worker and self.worker.is_alive():
            if not messagebox.askyesno("확인", "색인이 진행 중입니다. 중지하고 닫을까요?",
                                       parent=self):
                return
            self.stop_flag.set()
        self.destroy()


# ============================================================
# 메인 앱
# ============================================================
class FileFinderApp:
    def __init__(self, root):
        self.root = root
        self.root.title("파일 찾기")
        self.root.geometry("1200x840")
        self.root.minsize(820, 600)

        self.search_thread = None
        self.stop_flag = threading.Event()
        self.results = []
        self.index_mgr = IndexManager()
        self._preview_token = 0
        self._preview_image_ref = None  # Prevent GC

        self._build_ui()
        self.root.bind('<Escape>', lambda e: self._stop_search())
        self.root.bind('<Control-f>', lambda e: self.name_entry.focus_set())

    def _build_ui(self):
        style = ttk.Style()
        try:
            style.theme_use('vista')
        except tk.TclError:
            pass

        top = ttk.Frame(self.root, padding=(12, 10, 12, 6))
        top.pack(fill='x')

        ttk.Label(top, text="검색 폴더").grid(row=0, column=0, sticky='w', padx=(0, 8), pady=4)
        self.folder_var = tk.StringVar(value=str(Path.home()))
        self.folder_entry = ttk.Entry(top, textvariable=self.folder_var)
        self.folder_entry.grid(row=0, column=1, sticky='we', pady=4)
        ttk.Button(top, text="폴더 선택", command=self._browse, width=10).grid(row=0, column=2, padx=(6, 0))

        ttk.Label(top, text="파일명 포함").grid(row=1, column=0, sticky='w', padx=(0, 8), pady=4)
        self.name_var = tk.StringVar()
        self.name_entry = ttk.Entry(top, textvariable=self.name_var)
        self.name_entry.grid(row=1, column=1, sticky='we', pady=4)
        self.name_entry.bind('<Return>', lambda e: self._start_search())

        ttk.Label(top, text="확장자").grid(row=2, column=0, sticky='w', padx=(0, 8), pady=4)
        self.ext_var = tk.StringVar()
        ttk.Entry(top, textvariable=self.ext_var).grid(row=2, column=1, sticky='we', pady=4)
        ttk.Label(top, text="예: hwp, pdf, html  (쉼표로 구분)", foreground='#888').grid(row=2, column=2, sticky='w', padx=(6, 0))

        ttk.Label(top, text="파일 내용 포함").grid(row=3, column=0, sticky='w', padx=(0, 8), pady=4)
        self.content_var = tk.StringVar()
        content_entry = ttk.Entry(top, textvariable=self.content_var)
        content_entry.grid(row=3, column=1, sticky='we', pady=4)
        content_entry.bind('<Return>', lambda e: self._start_search())
        ttk.Label(top, text="지원: txt, html, py 등 텍스트 / hwp, pdf, docx, pptx, xlsx",
                  foreground='#888', font=('Segoe UI', 8)).grid(row=3, column=2, sticky='w', padx=(6, 0))

        opt = ttk.Frame(top)
        opt.grid(row=4, column=0, columnspan=3, sticky='w', pady=(6, 2))
        self.subdir_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(opt, text="하위 폴더 포함", variable=self.subdir_var).pack(side='left')
        self.case_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(opt, text="대소문자 구분", variable=self.case_var).pack(side='left', padx=(14, 0))
        self.skip_hidden_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(opt, text="숨김/시스템 폴더 제외", variable=self.skip_hidden_var).pack(side='left', padx=(14, 0))
        ttk.Label(opt, text="  |  내용 검색 최대 파일 크기 (MB)").pack(side='left', padx=(14, 2))
        self.maxsize_var = tk.StringVar(value="20")
        ttk.Entry(opt, textvariable=self.maxsize_var, width=5).pack(side='left')

        # 색인 옵션
        idx_opt = ttk.Frame(top)
        idx_opt.grid(row=5, column=0, columnspan=3, sticky='w', pady=(4, 0))
        self.use_index_var = tk.BooleanVar(value=False)
        idx_chk = ttk.Checkbutton(
            idx_opt, text="⚡ 색인으로 검색 (고속)",
            variable=self.use_index_var, command=self._on_index_toggle
        )
        idx_chk.pack(side='left')
        ttk.Button(idx_opt, text="색인 관리...", command=self._open_index_dialog, width=14).pack(side='left', padx=(10, 0))
        self.index_status = ttk.Label(idx_opt, text="", foreground='#444')
        self.index_status.pack(side='left', padx=(10, 0))
        self._update_index_status()

        btn = ttk.Frame(top)
        btn.grid(row=6, column=0, columnspan=3, pady=(8, 0))
        self.search_btn = ttk.Button(btn, text="🔍  검색 시작", command=self._start_search, width=14)
        self.search_btn.pack(side='left', padx=(0, 6))
        self.stop_btn = ttk.Button(btn, text="■  중지", command=self._stop_search, state='disabled', width=10)
        self.stop_btn.pack(side='left')
        ttk.Label(btn, text="  ·  Enter = 검색, Esc = 중지, 결과 더블클릭 = 탐색기에서 열기",
                  foreground='#888').pack(side='left', padx=(10, 0))

        top.columnconfigure(1, weight=1)

        ttk.Separator(self.root, orient='horizontal').pack(fill='x', padx=12, pady=(4, 0))

        # 결과 + 미리보기 — 수직 분할
        paned = ttk.PanedWindow(self.root, orient='vertical')
        paned.pack(fill='both', expand=True, padx=12, pady=(8, 6))

        mid = ttk.Frame(paned)
        paned.add(mid, weight=3)

        cols = ('name', 'path', 'size', 'modified', 'snippet')
        self.tree = ttk.Treeview(mid, columns=cols, show='headings', selectmode='extended')
        self.tree.heading('name', text='파일명', command=lambda: self._sort_by('name'))
        self.tree.heading('path', text='경로', command=lambda: self._sort_by('path'))
        self.tree.heading('size', text='크기', command=lambda: self._sort_by('size'))
        self.tree.heading('modified', text='수정일', command=lambda: self._sort_by('modified'))
        self.tree.heading('snippet', text='일치 내용')
        self.tree.column('name', width=220, anchor='w')
        self.tree.column('path', width=340, anchor='w')
        self.tree.column('size', width=80, anchor='e')
        self.tree.column('modified', width=130, anchor='w')
        self.tree.column('snippet', width=340, anchor='w')

        vsb = ttk.Scrollbar(mid, orient='vertical', command=self.tree.yview)
        self.tree.configure(yscrollcommand=vsb.set)
        self.tree.pack(side='left', fill='both', expand=True)
        vsb.pack(side='right', fill='y')

        self.tree.bind('<Double-Button-1>', self._open_in_explorer)
        self.tree.bind('<<TreeviewSelect>>', self._on_select)

        # 미리보기 패널
        preview = ttk.LabelFrame(paned, text=" 미리보기 ", padding=(8, 6, 8, 8))
        paned.add(preview, weight=2)

        header_frame = ttk.Frame(preview)
        header_frame.pack(fill='x')
        self.preview_header = ttk.Label(
            header_frame,
            text="결과를 선택하면 미리보기가 표시됩니다  ·  (텍스트 / HWP / PDF / DOCX / PPTX / XLSX / 이미지)",
            foreground='#888'
        )
        self.preview_header.pack(side='left')
        self.preview_meta = ttk.Label(header_frame, text="", foreground='#666')
        self.preview_meta.pack(side='right')

        ttk.Separator(preview, orient='horizontal').pack(fill='x', pady=(4, 6))

        # Text + Canvas 스택 (이미지일 때는 Canvas 사용)
        self.preview_container = ttk.Frame(preview)
        self.preview_container.pack(fill='both', expand=True)

        self.preview_text = tk.Text(
            self.preview_container, wrap='word', state='disabled',
            font=('맑은 고딕', 10), relief='flat',
            background='#FCFCFC', foreground='#111',
            padx=6, pady=4
        )
        self.preview_sb = ttk.Scrollbar(
            self.preview_container, orient='vertical',
            command=self.preview_text.yview
        )
        self.preview_text.configure(yscrollcommand=self.preview_sb.set)
        self.preview_text.pack(side='left', fill='both', expand=True)
        self.preview_sb.pack(side='right', fill='y')

        self.preview_canvas = tk.Canvas(
            self.preview_container, background='#1A1A1A',
            highlightthickness=0
        )
        # Canvas는 이미지 표시 시 pack

        # 하이라이트 태그
        self.preview_text.tag_configure(
            'hit', background='#FFE27A', foreground='#000'
        )
        self.preview_text.tag_configure(
            'dim', foreground='#888', font=('맑은 고딕', 9, 'italic')
        )

        self.menu = tk.Menu(self.root, tearoff=0)
        self.menu.add_command(label="탐색기에서 열기", command=self._open_in_explorer)
        self.menu.add_command(label="파일 실행", command=self._open_file)
        self.menu.add_command(label="경로 복사", command=self._copy_path)
        self.tree.bind('<Button-3>', self._show_context_menu)

        status = ttk.Frame(self.root, padding=(12, 4))
        status.pack(fill='x')
        self.status_var = tk.StringVar(value="준비 — 검색 폴더와 조건을 입력하세요")
        ttk.Label(status, textvariable=self.status_var).pack(side='left')
        self.count_var = tk.StringVar(value="")
        ttk.Label(status, textvariable=self.count_var, font=('Segoe UI', 9, 'bold')).pack(side='right')

    def _on_index_toggle(self):
        # 색인 모드에서는 "검색 폴더" 입력이 무시됨을 사용자에게 암시
        if self.use_index_var.get():
            self.folder_entry.state(['disabled'])
        else:
            self.folder_entry.state(['!disabled'])

    def _update_index_status(self):
        st = self.index_mgr.stats()
        if st['total_files'] == 0:
            self.index_status.configure(
                text="(색인 없음 — '색인 관리'에서 폴더를 추가하세요)",
                foreground='#888'
            )
        else:
            self.index_status.configure(
                text=f"색인: {st['total_files']:,}개 파일",
                foreground='#2C5F2D'
            )

    def _open_index_dialog(self):
        IndexDialog(self.root, self.index_mgr, on_close=self._update_index_status)

    def _browse(self):
        d = filedialog.askdirectory(initialdir=self.folder_var.get() or str(Path.home()))
        if d:
            self.folder_var.set(os.path.normpath(d))

    def _start_search(self):
        if self.search_thread and self.search_thread.is_alive():
            return

        use_index = self.use_index_var.get()

        if use_index:
            st = self.index_mgr.stats()
            if st['total_files'] == 0:
                messagebox.showinfo(
                    "알림",
                    "아직 색인된 파일이 없습니다.\n'색인 관리'에서 폴더를 추가하고 색인을 만들어주세요."
                )
                return
        else:
            folder = self.folder_var.get().strip()
            if not folder or not os.path.isdir(folder):
                messagebox.showerror("오류", "검색 폴더가 존재하지 않습니다.")
                return

        if (not self.name_var.get().strip()
                and not self.ext_var.get().strip()
                and not self.content_var.get().strip()):
            if not messagebox.askyesno("확인", "검색어가 비어 있습니다. 모든 파일을 나열할까요?"):
                return

        if not use_index and self.content_var.get().strip() and not self.ext_var.get().strip():
            if not messagebox.askyesno(
                "권장 설정",
                "파일 내용 검색은 속도가 느립니다. 확장자를 지정하거나 색인 검색 모드를 권장합니다.\n\n"
                "그래도 계속할까요?"
            ):
                return

        try:
            max_mb = float(self.maxsize_var.get())
        except ValueError:
            messagebox.showerror("오류", "최대 파일 크기는 숫자로 입력하세요.")
            return

        for i in self.tree.get_children():
            self.tree.delete(i)
        self.results = []

        self.stop_flag.clear()
        self.search_btn.configure(state='disabled')
        self.stop_btn.configure(state='normal')
        self.status_var.set("검색 중...")
        self.count_var.set("")

        if use_index:
            self.search_thread = threading.Thread(
                target=self._index_search_worker, daemon=True
            )
        else:
            self.search_thread = threading.Thread(
                target=self._live_search_worker, args=(max_mb,), daemon=True
            )
        self.search_thread.start()

    def _stop_search(self):
        if self.search_thread and self.search_thread.is_alive():
            self.stop_flag.set()

    def _index_search_worker(self):
        name_q = self.name_var.get().strip()
        ext_q = self.ext_var.get().strip()
        content_q = self.content_var.get().strip()
        case_sensitive = self.case_var.get()
        try:
            t0 = time.time()
            rows = self.index_mgr.search(name_q, ext_q, content_q, case_sensitive)
            elapsed = time.time() - t0
            for path, name, size, mtime, snippet in rows:
                if self.stop_flag.is_set():
                    break
                root_dir = os.path.dirname(path)
                self.results.append((name, root_dir, size or 0, mtime or 0, snippet or ''))
            self.root.after(0, self._finish_index_search, elapsed, len(rows))
        except Exception as e:
            self.root.after(0, lambda: messagebox.showerror("오류", f"색인 검색 오류: {e}"))
            self.root.after(0, self._finish_index_search, 0, 0)

    def _finish_index_search(self, elapsed, count):
        self._flush_results()
        self.search_btn.configure(state='normal')
        self.stop_btn.configure(state='disabled')
        self.status_var.set(f"완료 (색인 검색 · {elapsed*1000:,.0f} ms)")
        self.count_var.set(f"{count:,}건 발견")

    def _live_search_worker(self, max_mb):
        folder = self.folder_var.get()
        name_q = self.name_var.get().strip()
        ext_q = self.ext_var.get().strip()
        content_q = self.content_var.get().strip()
        subdir = self.subdir_var.get()
        case_sensitive = self.case_var.get()
        skip_hidden = self.skip_hidden_var.get()
        max_bytes = int(max_mb * 1024 * 1024)

        name_q_cmp = name_q if case_sensitive else name_q.lower()
        exts = None
        if ext_q:
            exts = set(e.strip().lstrip('.').lower() for e in ext_q.split(',') if e.strip())

        count = 0
        scanned = 0
        content_scanned = 0
        try:
            if subdir:
                walker = os.walk(folder, onerror=lambda e: None)
            else:
                try:
                    entries = os.listdir(folder)
                except OSError:
                    entries = []
                files_only = [e for e in entries if os.path.isfile(os.path.join(folder, e))]
                walker = [(folder, [], files_only)]

            for root_dir, dirs, files in walker:
                if self.stop_flag.is_set():
                    break
                if skip_hidden:
                    dirs[:] = [d for d in dirs if not self._is_hidden(os.path.join(root_dir, d))]
                for f in files:
                    if self.stop_flag.is_set():
                        break
                    scanned += 1
                    fname_cmp = f if case_sensitive else f.lower()
                    if name_q and name_q_cmp not in fname_cmp:
                        continue
                    if exts is not None:
                        fext = f.rsplit('.', 1)[-1].lower() if '.' in f else ''
                        if fext not in exts:
                            continue
                    fpath = os.path.join(root_dir, f)
                    try:
                        st = os.stat(fpath)
                    except OSError:
                        continue
                    snippet = ''
                    if content_q:
                        content_scanned += 1
                        if content_scanned % 5 == 0 or st.st_size > 1_000_000:
                            self.root.after(0, self._status_now,
                                            f"내용 검색 중 · {f[:60]}", scanned, count)
                        text = extract_text_from_file(fpath, max_bytes)
                        if text is None:
                            continue
                        snippet = find_snippet(text, content_q, case_sensitive)
                        if not snippet:
                            continue
                    self.results.append((f, root_dir, st.st_size, st.st_mtime, snippet))
                    count += 1
                    flush_every = 10 if content_q else 30
                    if count % flush_every == 0:
                        self.root.after(0, self._flush_and_status, scanned, count)
                if scanned % 800 == 0 and not content_q:
                    self.root.after(0, self._status_only, scanned, count)
        except Exception as e:
            self.root.after(0, lambda: messagebox.showerror("오류", f"검색 중 오류: {e}"))
        self.root.after(0, self._finish_live_search, scanned, count)

    @staticmethod
    def _is_hidden(path):
        name = os.path.basename(path)
        if name.startswith('.'):
            return True
        try:
            import ctypes
            attrs = ctypes.windll.kernel32.GetFileAttributesW(str(path))
            if attrs == -1:
                return False
            return bool(attrs & 2) or bool(attrs & 4)
        except Exception:
            return False

    def _flush_and_status(self, scanned, count):
        self._flush_results()
        self.status_var.set(f"검색 중... {scanned:,} 파일 스캔")
        self.count_var.set(f"{count:,}건 발견")

    def _status_only(self, scanned, count):
        self.status_var.set(f"검색 중... {scanned:,} 파일 스캔")
        self.count_var.set(f"{count:,}건 발견")

    def _status_now(self, msg, scanned, count):
        self._flush_results()
        self.status_var.set(f"{msg}  ({scanned:,} 스캔)")
        self.count_var.set(f"{count:,}건 발견")

    def _flush_results(self):
        shown = len(self.tree.get_children())
        for item in self.results[shown:]:
            f, root_dir, size, mtime, snippet = item
            size_str = self._fmt_size(size) if size else ''
            mtime_str = (datetime.datetime.fromtimestamp(mtime).strftime('%Y-%m-%d %H:%M')
                         if mtime else '')
            self.tree.insert('', 'end', values=(f, root_dir, size_str, mtime_str, snippet or ''))

    def _finish_live_search(self, scanned, count):
        self._flush_results()
        self.search_btn.configure(state='normal')
        self.stop_btn.configure(state='disabled')
        if self.stop_flag.is_set():
            self.status_var.set(f"중지됨 — {scanned:,} 파일 스캔")
        else:
            self.status_var.set(f"완료 — {scanned:,} 파일 스캔")
        self.count_var.set(f"{count:,}건 발견")

    @staticmethod
    def _fmt_size(size):
        for unit in ('B', 'KB', 'MB', 'GB'):
            if size < 1024:
                if unit == 'B':
                    return f"{size:,.0f} B"
                return f"{size:,.1f} {unit}"
            size /= 1024
        return f"{size:,.1f} TB"

    def _sort_by(self, col):
        items = [(self.tree.set(k, col), k) for k in self.tree.get_children('')]
        def size_key(v):
            try:
                n, u = v.rsplit(' ', 1)
                m = {'B': 1, 'KB': 1024, 'MB': 1024**2, 'GB': 1024**3, 'TB': 1024**4}.get(u, 1)
                return float(n.replace(',', '')) * m
            except Exception:
                return 0.0
        if col == 'size':
            items.sort(key=lambda x: size_key(x[0]))
        else:
            items.sort(key=lambda x: x[0])
        if getattr(self, '_last_sort', None) == col:
            items.reverse()
            self._last_sort = None
        else:
            self._last_sort = col
        for idx, (_, k) in enumerate(items):
            self.tree.move(k, '', idx)

    def _selected_fullpath(self):
        sel = self.tree.selection()
        if not sel:
            return None
        vals = self.tree.item(sel[0], 'values')
        if not vals:
            return None
        return os.path.join(vals[1], vals[0])

    def _open_in_explorer(self, event=None):
        fp = self._selected_fullpath()
        if event and not fp:
            row = self.tree.identify_row(event.y)
            if row:
                self.tree.selection_set(row)
                fp = self._selected_fullpath()
        if not fp:
            return
        try:
            subprocess.Popen(['explorer', '/select,', fp])
        except Exception as e:
            messagebox.showerror("오류", f"탐색기 열기 실패: {e}")

    def _open_file(self):
        fp = self._selected_fullpath()
        if not fp:
            return
        try:
            os.startfile(fp)
        except Exception as e:
            messagebox.showerror("오류", f"파일 열기 실패: {e}")

    def _copy_path(self):
        fp = self._selected_fullpath()
        if not fp:
            return
        self.root.clipboard_clear()
        self.root.clipboard_append(fp)
        self.status_var.set(f"경로 복사됨: {fp}")

    def _show_context_menu(self, event):
        row = self.tree.identify_row(event.y)
        if row:
            self.tree.selection_set(row)
            try:
                self.menu.tk_popup(event.x_root, event.y_root)
            finally:
                self.menu.grab_release()

    # ===== 미리보기 =====
    IMAGE_EXTS = {'png', 'gif', 'jpg', 'jpeg', 'bmp', 'webp', 'tiff', 'tif', 'ico'}
    PREVIEW_MAX_CHARS = 50_000
    PREVIEW_MAX_FILE_MB = 50

    def _on_select(self, event=None):
        sel = self.tree.selection()
        if not sel:
            return
        vals = self.tree.item(sel[0], 'values')
        if not vals:
            return
        fp = os.path.join(vals[1], vals[0])

        self._preview_token += 1
        token = self._preview_token

        fname = os.path.basename(fp)
        self.preview_header.configure(text=f"{fname}  ·  로딩 중...", foreground='#444')
        self.preview_meta.configure(text='')
        self._show_text_mode()
        self._set_preview_text('', [])

        threading.Thread(
            target=self._load_preview, args=(fp, token), daemon=True
        ).start()

    def _load_preview(self, fp, token):
        try:
            st = os.stat(fp)
        except OSError as e:
            self.root.after(0, self._preview_error, fp, f"파일 읽기 실패: {e}", token)
            return

        ext = fp.rsplit('.', 1)[-1].lower() if '.' in fp else ''
        size_mb = st.st_size / (1024 * 1024)

        # 이미지
        if ext in self.IMAGE_EXTS:
            if size_mb > 30:
                self.root.after(0, self._preview_error, fp,
                                f"이미지가 큽니다 ({size_mb:.1f} MB) — 미리보기 생략", token)
                return
            try:
                from PIL import Image
                img = Image.open(fp)
                img.load()
                # 메타데이터
                meta = f"{img.format} · {img.width}×{img.height} · {self._fmt_size(st.st_size)}"
                if token != self._preview_token:
                    return
                self.root.after(0, self._preview_show_image, fp, img, meta, token)
            except Exception as e:
                self.root.after(0, self._preview_error, fp, f"이미지 로드 실패: {e}", token)
            return

        # 텍스트/문서
        max_bytes = int(self.PREVIEW_MAX_FILE_MB * 1024 * 1024)
        if size_mb > self.PREVIEW_MAX_FILE_MB:
            self.root.after(
                0, self._preview_error, fp,
                f"파일이 너무 큽니다 ({size_mb:.1f} MB > {self.PREVIEW_MAX_FILE_MB} MB) — 미리보기 생략",
                token
            )
            return

        # 1단계: 빠른 미리보기 (앞 페이지/섹션만)
        text = extract_preview_text(fp, max_bytes)
        is_partial = True

        # PDF: 텍스트가 너무 적으면 스캔 PDF로 간주 → 이미지 렌더링
        if ext == 'pdf' and (text is None or len(text.strip()) < 50):
            img = render_pdf_first_page_as_image(fp)
            if img is not None:
                meta = f"PDF · 스캔/이미지 기반 · {self._fmt_size(st.st_size)}  ·  첫 페이지"
                if token != self._preview_token:
                    return
                self.root.after(0, self._preview_show_image, fp, img, meta, token)
                return

        # 검색어 목록
        queries = []
        content_q = self.content_var.get().strip()
        name_q = self.name_var.get().strip()
        for q in (content_q, name_q):
            if q:
                queries.append(q)

        case_sensitive = self.case_var.get()

        # 2단계: 검색어가 있는데 빠른 미리보기에 안 보이면 전체 추출
        if queries and text is not None:
            haystack = text if case_sensitive else text.lower()
            missing = [q for q in queries
                       if (q if case_sensitive else q.lower()) not in haystack]
            if missing:
                if token != self._preview_token:
                    return
                self.root.after(
                    0, lambda: self.preview_meta.configure(
                        text="검색어 확인을 위해 전체 내용 로드 중..."
                    )
                )
                full_text = extract_text_from_file(fp, max_bytes)
                if full_text is not None:
                    text = full_text
                    is_partial = False

        # 빠른 미리보기가 None이면 전체라도 시도
        if text is None:
            text = extract_text_from_file(fp, max_bytes)
            is_partial = False

        if text is None:
            self.root.after(
                0, self._preview_error, fp,
                "(지원하지 않는 형식이거나 텍스트 추출 실패)",
                token
            )
            return

        truncated = False
        if len(text) > self.PREVIEW_MAX_CHARS:
            text = text[:self.PREVIEW_MAX_CHARS]
            truncated = True

        # 검색어 위치
        hits = []
        if queries:
            haystack = text if case_sensitive else text.lower()
            for q in queries:
                needle = q if case_sensitive else q.lower()
                start = 0
                while True:
                    idx = haystack.find(needle, start)
                    if idx < 0:
                        break
                    hits.append((idx, idx + len(q)))
                    start = idx + len(q)
                    if len(hits) > 500:
                        break

        if token != self._preview_token:
            return
        meta = f"{self._fmt_size(st.st_size)}  ·  {datetime.datetime.fromtimestamp(st.st_mtime).strftime('%Y-%m-%d %H:%M')}"
        if is_partial:
            meta += "  ·  빠른 미리보기"
        if truncated:
            meta += "  ·  5만자 이후 생략"
        self.root.after(0, self._preview_show_text, fp, text, hits, meta, token)

    def _show_text_mode(self):
        self.preview_canvas.pack_forget()
        if not self.preview_text.winfo_ismapped():
            self.preview_text.pack(side='left', fill='both', expand=True)
            self.preview_sb.pack(side='right', fill='y')

    def _show_image_mode(self):
        self.preview_text.pack_forget()
        self.preview_sb.pack_forget()
        if not self.preview_canvas.winfo_ismapped():
            self.preview_canvas.pack(fill='both', expand=True)

    def _preview_show_text(self, fp, text, hits, meta, token):
        if token != self._preview_token:
            return
        self._show_text_mode()
        self.preview_header.configure(text=os.path.basename(fp), foreground='#111')
        self.preview_meta.configure(text=meta)
        self._set_preview_text(text, hits)

    def _set_preview_text(self, text, hits):
        self.preview_text.configure(state='normal')
        self.preview_text.delete('1.0', 'end')
        self.preview_text.insert('1.0', text)
        # 하이라이트
        if hits and text:
            for start, end in hits:
                start_idx = f"1.0+{start}c"
                end_idx = f"1.0+{end}c"
                self.preview_text.tag_add('hit', start_idx, end_idx)
            # 첫 매치로 스크롤
            self.preview_text.see(f"1.0+{hits[0][0]}c")
        self.preview_text.configure(state='disabled')

    def _preview_show_image(self, fp, img, meta, token):
        if token != self._preview_token:
            return
        self._show_image_mode()
        self.preview_header.configure(text=os.path.basename(fp), foreground='#111')
        self.preview_meta.configure(text=meta)
        self._render_image_fitted(img)
        # 창 크기 변경 시 재조정
        self.preview_canvas.bind(
            '<Configure>',
            lambda e, i=img: self._render_image_fitted(i)
        )

    def _render_image_fitted(self, img):
        try:
            from PIL import ImageTk
            cw = max(self.preview_canvas.winfo_width(), 100)
            ch = max(self.preview_canvas.winfo_height(), 100)
            # 비율 유지
            iw, ih = img.size
            scale = min(cw / iw, ch / ih, 1.0)
            nw, nh = max(1, int(iw * scale)), max(1, int(ih * scale))
            resized = img.resize((nw, nh))
            photo = ImageTk.PhotoImage(resized)
            self._preview_image_ref = photo  # keep reference
            self.preview_canvas.delete('all')
            self.preview_canvas.create_image(cw // 2, ch // 2, anchor='center', image=photo)
        except Exception as e:
            self.preview_canvas.delete('all')
            self.preview_canvas.create_text(
                50, 50, anchor='nw',
                text=f"이미지 표시 실패: {e}", fill='#FFF', font=('맑은 고딕', 10)
            )

    def _preview_error(self, fp, msg, token):
        if token != self._preview_token:
            return
        self._show_text_mode()
        try:
            st = os.stat(fp)
            meta = f"{self._fmt_size(st.st_size)}  ·  {datetime.datetime.fromtimestamp(st.st_mtime).strftime('%Y-%m-%d %H:%M')}"
        except Exception:
            meta = ''
        self.preview_header.configure(text=os.path.basename(fp), foreground='#111')
        self.preview_meta.configure(text=meta)
        self.preview_text.configure(state='normal')
        self.preview_text.delete('1.0', 'end')
        self.preview_text.insert('1.0', msg, 'dim')
        self.preview_text.configure(state='disabled')


def main():
    root = tk.Tk()
    FileFinderApp(root)
    root.mainloop()


if __name__ == '__main__':
    main()
